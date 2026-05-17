"""Generate Insider Threat Commit Intelligence client presentation PPT."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colours ──
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(239, 68, 68)
ORANGE = RGBColor(249, 115, 22)
YELLOW = RGBColor(234, 179, 8)
BLUE = RGBColor(99, 102, 241)
CYAN = RGBColor(6, 182, 212)
GREEN = RGBColor(16, 185, 129)
GREY = RGBColor(148, 163, 184)
LIGHT_GREY = RGBColor(203, 213, 225)
DIM = RGBColor(100, 116, 139)
DARK_BG = RGBColor(15, 23, 42)
CARD_BG = RGBColor(20, 30, 50)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helpers ──
def set_slide_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def add_text(tf, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_after=Pt(6)):
    p = tf.paragraphs[-1] if tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = space_after
    return p


def add_rect(slide, left, top, width, height, fill_color=CARD_BG, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def title_slide_header(slide, section_num, title):
    """Add section number + title bar at top."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.1), fill_color=DARK_BG, border_color=None)
    tb = add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{section_num}  "
    run1.font.size = Pt(14)
    run1.font.color.rgb = BLUE
    run1.font.bold = True
    run2 = p.add_run()
    run2.text = title
    run2.font.size = Pt(28)
    run2.font.color.rgb = WHITE
    run2.font.bold = True


def bullet_list(tf, items, size=16, color=LIGHT_GREY, bullet_color=BLUE):
    for item in items:
        p = tf.add_paragraph()
        p.space_after = Pt(4)
        p.level = 0
        # bullet char
        run_b = p.add_run()
        run_b.text = "  \u2022  "
        run_b.font.size = Pt(size)
        run_b.font.color.rgb = bullet_color
        run_b.font.bold = True
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(size)
        run_t.font.color.rgb = color


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Accent line
add_rect(slide, Inches(0.8), Inches(1.8), Inches(0.6), Pt(4), fill_color=RED)

tb = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(10), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Detecting Insider Threats\nThrough Commit Intelligence"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.line_spacing = Pt(54)

tb2 = add_textbox(slide, Inches(0.8), Inches(3.8), Inches(9), Inches(0.8))
tf2 = tb2.text_frame
tf2.word_wrap = True
add_text(tf2, "Every code commit contains rich behavioral and technical signals. By analyzing metadata,\nfile changes, code diffs, and developer patterns, we detect insider threat indicators in real-time.", size=16, color=GREY)

# Stats row
stats = [("13", "Threat Signals", RED), ("4", "Data Layers", ORANGE), ("58+", "Detection Rules", BLUE), ("0\u201310", "Risk Scoring", CYAN)]
for i, (num, label, color) in enumerate(stats):
    x = Inches(0.8 + i * 2.6)
    y = Inches(5.2)
    tb = add_textbox(slide, x, y, Inches(2.2), Inches(1.2))
    tf = tb.text_frame
    add_text(tf, num, size=40, bold=True, color=color, align=PP_ALIGN.LEFT)
    add_text(tf, label, size=14, color=DIM, align=PP_ALIGN.LEFT)

# Footer badge
tb = add_textbox(slide, Inches(0.8), Inches(6.7), Inches(4), Inches(0.4))
tf = tb.text_frame
add_text(tf, "INSIDER THREAT INTELLIGENCE  |  SECUREDEV AI", size=10, color=DIM, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — How It Works
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "01", "How It Works")

steps = [
    ("1", "Commit Pushed", "Developer pushes code\nto GitHub / GitLab"),
    ("2", "Data Extracted", "API fetches metadata,\nstats, files, and diff"),
    ("3", "Multi-Layer Scan", "13 signal detectors analyze\neach data layer"),
    ("4", "Baseline Check", "Compare against developer's\nbehavioral baseline"),
    ("5", "Risk Score", "Weighted 0-10 score with\nrisk level and alert"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    y = Inches(2.0)
    card = add_rect(slide, x, y, Inches(2.2), Inches(3.5), fill_color=CARD_BG, border_color=RGBColor(40, 50, 70))
    tb = add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(1.8), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    add_text(tf, num, size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER, space_after=Pt(8))
    add_text(tf, title, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=Pt(10))
    add_text(tf, desc, size=13, color=GREY, align=PP_ALIGN.CENTER)
    # arrow
    if i < 4:
        tb_arr = add_textbox(slide, x + Inches(2.2), y + Inches(1.3), Inches(0.3), Inches(0.5))
        tfa = tb_arr.text_frame
        add_text(tfa, "\u2192", size=24, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Anatomy of a Commit (4 Data Layers)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "02", "Anatomy of a Git Commit \u2014 4 Data Layers")

layers = [
    ("Layer 1", "Metadata \u2014 The Envelope", BLUE,
     ["author.name / author.email", "committer.name / committer.email",
      "author.date (timestamp)", "verification (GPG signed?)", "message (description)"]),
    ("Layer 2", "Statistics \u2014 The Size", ORANGE,
     ["stats.additions (lines added)", "stats.deletions (lines deleted)",
      "stats.total (combined)", "files_changed (count)"]),
    ("Layer 3", "Files Changed \u2014 The Surface", GREEN,
     ["filename (full path)", "status (added/modified/removed)",
      "per-file additions & deletions", "previous_filename (if renamed)"]),
    ("Layer 4", "The Diff \u2014 Actual Code", RED,
     ["Line-by-line unified diff", "+ lines = code added",
      "\u2013 lines = code removed", "Scanned against 58+ SAST rules"]),
]
for i, (badge, title, color, fields) in enumerate(layers):
    x = Inches(0.4 + i * 3.15)
    y = Inches(1.5)
    card = add_rect(slide, x, y, Inches(3.0), Inches(5.2), fill_color=CARD_BG, border_color=color)
    tb = add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.6), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    add_text(tf, badge, size=11, bold=True, color=color, space_after=Pt(4))
    add_text(tf, title, size=17, bold=True, color=WHITE, space_after=Pt(12))
    for f in fields:
        p = tf.add_paragraph()
        p.space_after = Pt(6)
        run_b = p.add_run()
        run_b.text = "\u2022  "
        run_b.font.size = Pt(13)
        run_b.font.color.rgb = color
        run_t = p.add_run()
        run_t.text = f
        run_t.font.size = Pt(13)
        run_t.font.color.rgb = LIGHT_GREY


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — 13 Threat Signals Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "03", "13 Insider Threat Signals We Detect")

signals_overview = [
    ("1", "Malicious Code Patterns (SAST)", "+2.5", "Critical", RED),
    ("2", "Secrets & Credentials in Code", "+2.5", "Critical", RED),
    ("3", "Sensitive File Access", "+1.5/file", "High", ORANGE),
    ("4", "Author / Committer Mismatch", "+1.5", "High", ORANGE),
    ("5", "Off-Hours Commit Activity", "+1.0", "Medium", YELLOW),
    ("6", "Large Code Deletions", "+1.0", "Medium", YELLOW),
    ("7", "Unsigned Commits", "+0.5", "Medium", YELLOW),
    ("8", "Force Push (History Rewrite)", "+2.0", "Critical", RED),
    ("9", "Suspicious Commit Messages", "+0.5", "Medium", YELLOW),
    ("10", "Binary File Injection", "+1.0/file", "High", ORANGE),
    ("11", "Dependency Manipulation", "+1.5", "High", ORANGE),
    ("12", "CI/CD Pipeline Tampering", "+1.5/find", "Critical", RED),
    ("13", "Config & Permission Weakening", "+1.0/find", "High", ORANGE),
]

# Two-column layout
for idx, (num, name, pts, sev, color) in enumerate(signals_overview):
    col = 0 if idx < 7 else 1
    row = idx if idx < 7 else idx - 7
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.5 + row * 0.78)
    card = add_rect(slide, x, y, Inches(6.0), Inches(0.65), fill_color=CARD_BG, border_color=RGBColor(40, 50, 70))
    # Number
    tb = add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.5), Inches(0.5))
    tf = tb.text_frame
    add_text(tf, num, size=16, bold=True, color=color)
    # Name
    tb = add_textbox(slide, x + Inches(0.6), y + Inches(0.08), Inches(3.6), Inches(0.5))
    tf = tb.text_frame
    add_text(tf, name, size=14, bold=False, color=WHITE)
    # Points
    tb = add_textbox(slide, x + Inches(4.2), y + Inches(0.08), Inches(0.8), Inches(0.5))
    tf = tb.text_frame
    add_text(tf, pts, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Severity badge
    tb = add_textbox(slide, x + Inches(5.0), y + Inches(0.08), Inches(0.9), Inches(0.5))
    tf = tb.text_frame
    add_text(tf, sev, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Signal Deep Dive: Critical Signals
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "03a", "Critical Signals \u2014 Deep Dive")

# Left: SAST
add_rect(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.5), fill_color=CARD_BG, border_color=RED)
tb = add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True
add_text(tf, "1. Malicious Code Patterns (SAST on Diffs)  |  +2.5 pts", size=16, bold=True, color=RED, space_after=Pt(10))
add_text(tf, "Every added line scanned against 58 insider threat-specific rules:", size=13, color=GREY, space_after=Pt(8))
categories = [
    "Backdoor Insertion \u2014 hardcoded bypass credentials, hidden admin routes",
    "Data Exfiltration \u2014 outbound HTTP to external endpoints, base64 data piping",
    "Logging Tampering \u2014 disabling security logs, redirecting audit trails",
    "Access Control Weakening \u2014 removing auth checks, broadening permissions",
    "Obfuscation \u2014 base64/hex encoded payloads, eval() with dynamic strings",
    "Crypto Weakening \u2014 downgrading encryption, weak ciphers, hardcoded keys",
]
bullet_list(tf, categories, size=12, color=LIGHT_GREY, bullet_color=RED)
add_text(tf, "", size=6, color=BLACK, space_after=Pt(4))
add_text(tf, "Example: Backdoor Detected", size=12, bold=True, color=ORANGE, space_after=Pt(4))
add_text(tf, '+ if user == "svc_maint" and pwd == "x9K#backdoor":', size=11, color=GREEN, space_after=Pt(2))
add_text(tf, '+     return AuthResult(authenticated=True, role="admin")', size=11, color=GREEN)

# Right: Force Push + Secrets
add_rect(slide, Inches(6.8), Inches(1.4), Inches(6.2), Inches(2.6), fill_color=CARD_BG, border_color=RED)
tb = add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(2.4))
tf = tb.text_frame
tf.word_wrap = True
add_text(tf, "8. Force Push (History Rewrite)  |  +2.0 pts", size=16, bold=True, color=RED, space_after=Pt(8))
add_text(tf, "Replaces existing commits \u2014 erases evidence permanently.", size=13, color=GREY, space_after=Pt(6))
bullet_list(tf, [
    "Push malicious code \u2192 wait for deploy \u2192 force push to erase",
    "Remove audit logs that were previously committed",
    "Rewrite commit authorship to blame others",
], size=12, color=LIGHT_GREY, bullet_color=RED)

add_rect(slide, Inches(6.8), Inches(4.2), Inches(6.2), Inches(2.7), fill_color=CARD_BG, border_color=RED)
tb = add_textbox(slide, Inches(7.0), Inches(4.4), Inches(5.8), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
add_text(tf, "2. Secrets & Credentials in Code  |  +2.5 pts", size=16, bold=True, color=RED, space_after=Pt(8))
add_text(tf, "API keys, tokens, passwords, private keys in commits.", size=13, color=GREY, space_after=Pt(6))
bullet_list(tf, [
    "AWS Keys: AKIA[0-9A-Z]{16}",
    "GitHub Tokens: ghp_[a-zA-Z0-9]{36}",
    "Private Keys: -----BEGIN RSA PRIVATE KEY-----",
    "Generic passwords: password = \"hardcoded_value\"",
], size=12, color=LIGHT_GREY, bullet_color=RED)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Signal Deep Dive: High & Medium Signals
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "03b", "High & Medium Signals \u2014 Deep Dive")

# 2x3 grid
cards = [
    ("3. Sensitive File Access", "+1.5/file", ORANGE, [
        ".env, .env.production \u2014 database passwords, API keys",
        "id_rsa, *.pem, *.pfx \u2014 private keys & certificates",
        "terraform.tfstate \u2014 cloud credentials in plaintext",
        "shadow, passwd \u2014 system auth databases",
    ]),
    ("4. Author/Committer Mismatch", "+1.5", ORANGE, [
        "author.email \u2260 committer.email",
        "Credential sharing between developers",
        "Impersonation \u2014 attacker attributes to trusted dev",
        "Compromised account in different timezone",
    ]),
    ("5. Off-Hours Commits", "+1.0", YELLOW, [
        "Personalized per developer (not fixed hours)",
        "Rolling 60-commit baseline, z-score analysis",
        "z \u2265 2.0 \u2192 Medium  |  z \u2265 3.0 \u2192 High anomaly",
        "3 AM commit from 9-5 dev = z-score 4.5",
    ]),
    ("6. Large Deletions", "+1.0", YELLOW, [
        "Absolute: > 500 lines deleted",
        "Behavioral: > 5x developer's baseline average",
        "Sabotage, cover tracks, IP theft preparation",
        "Denial of service via migration file deletion",
    ]),
    ("7. Unsigned Commits", "+0.5", YELLOW, [
        "Missing GPG/SSH signature verification",
        "Attacker has PAT but not signing key",
        "Supporting signal in compound detection",
        "Unsigned + off-hours + sensitive file = strong",
    ]),
    ("9. Suspicious Messages", "+0.5", YELLOW, [
        "Only on non-trivial commits (>50 lines)",
        'Empty or <4 chars: "", "."',
        'Single-word vague: "fix", "wip", "stuff"',
        "487 lines added with message '.' \u2192 flagged",
    ]),
]
for i, (title, pts, color, items) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = Inches(0.3 + col * 4.3)
    y = Inches(1.4 + row * 3.0)
    card = add_rect(slide, x, y, Inches(4.1), Inches(2.8), fill_color=CARD_BG, border_color=color)
    tb = add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(3.7), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = title
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = f"  {pts}"
    r2.font.size = Pt(14)
    r2.font.bold = True
    r2.font.color.rgb = color
    p.space_after = Pt(6)
    bullet_list(tf, items, size=11, color=LIGHT_GREY, bullet_color=color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — New Signals: Supply Chain & Infrastructure
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "03c", "Supply Chain & Infrastructure Signals")

# Binary + Dependency + CI/CD + Config
new_cards = [
    ("10. Binary File Injection", "+1.0/file (cap 2.0)", ORANGE, [
        "Executables: .exe .dll .so .dylib .bin .msi",
        "Archives: .zip .tar .rar .7z",
        "Database dumps: .sql .sqlite .dump .bak",
        "Binary files bypass code review \u2014 can't diff them",
        "Malware delivery, data exfiltration via .zip",
    ]),
    ("11. Dependency Manipulation", "+1.5", ORANGE, [
        "Typosquat suspect (l0dash, requ3sts)",
        "Vulnerable version pinning (log4j 2.14)",
        "Custom registry URL \u2014 supply chain hijack",
        "preinstall/postinstall scripts on npm install",
        "Removing security deps (helmet, bcrypt, bandit)",
    ]),
    ("12. CI/CD Pipeline Tampering", "+1.5/find (cap 3.0)", RED, [
        "Security scan disabled (Snyk/SonarQube/Trivy)",
        "Hook bypass: --no-verify flags",
        "Secret exfiltration: echo $SECRET | curl",
        "Deploy target changes without PR flow",
        "Untrusted base images in Dockerfile",
    ]),
    ("13. Config & Permission Weakening", "+1.0/find (cap 2.0)", ORANGE, [
        "CORS wildcard (*) \u2014 any site can call your API",
        "MFA disabled, auth requirement removed",
        "Rate limit disabled, debug mode enabled",
        ".gitignore unignored \u2014 secrets get tracked",
        "CODEOWNERS removed \u2014 bypasses review gates",
    ]),
]
for i, (title, pts, color, items) in enumerate(new_cards):
    col = i % 2
    row = i // 2
    x = Inches(0.3 + col * 6.5)
    y = Inches(1.4 + row * 3.0)
    card = add_rect(slide, x, y, Inches(6.2), Inches(2.8), fill_color=CARD_BG, border_color=color)
    tb = add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.8), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = title
    r1.font.size = Pt(15)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = f"  |  {pts}"
    r2.font.size = Pt(13)
    r2.font.bold = True
    r2.font.color.rgb = color
    p.space_after = Pt(8)
    bullet_list(tf, items, size=12, color=LIGHT_GREY, bullet_color=color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Behavioral Baseline Engine
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "04", "Behavioral Baseline Engine")

# Left card — metrics
add_rect(slide, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.5), fill_color=CARD_BG, border_color=BLUE)
tb = add_textbox(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True
add_text(tf, "Per-Developer Behavioral Profiling", size=18, bold=True, color=WHITE, space_after=Pt(6))
add_text(tf, "Rolling window of 60 commits per developer", size=13, color=GREY, space_after=Pt(12))
add_text(tf, "Baseline Metrics Computed:", size=14, bold=True, color=BLUE, space_after=Pt(8))
metrics = [
    "Mean Commit Hour \u2014 average time of day",
    "Std Dev (Hours) \u2014 how much time varies",
    "Avg Additions / Deletions \u2014 typical change size",
    "P90 Additions \u2014 their \"large\" commit threshold",
    "Avg Files Changed \u2014 typical scope",
    "Avg Risk Score \u2014 historical risk level",
    "Commits/Week \u2014 activity frequency",
]
bullet_list(tf, metrics, size=12, color=LIGHT_GREY, bullet_color=BLUE)

# Right card — maturity
add_rect(slide, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.5), fill_color=CARD_BG, border_color=BLUE)
tb = add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True
add_text(tf, "Baseline Maturity Progression", size=18, bold=True, color=WHITE, space_after=Pt(12))

add_text(tf, "< 5 commits \u2192 Insufficient", size=15, bold=True, color=RED, space_after=Pt(4))
add_text(tf, "Not enough data. No anomaly detection active.\nMonitoring only.", size=12, color=GREY, space_after=Pt(14))

add_text(tf, "5\u201319 commits \u2192 Partial", size=15, bold=True, color=YELLOW, space_after=Pt(4))
add_text(tf, "Basic window-based checks (outside typical hours).\nNo z-score analysis yet.", size=12, color=GREY, space_after=Pt(14))

add_text(tf, "20+ commits \u2192 Established", size=15, bold=True, color=GREEN, space_after=Pt(4))
add_text(tf, "Full statistical analysis: z-scores for timing,\nmultiplier-based size detection, risk trend analysis.\nHigh confidence anomaly detection.", size=12, color=GREY, space_after=Pt(8))

add_text(tf, "Key insight: A commit that's normal for a senior\nbackend engineer may be highly anomalous for\na frontend intern.", size=13, bold=False, color=CYAN, space_after=Pt(4))


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Risk Scoring Model
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "05", "Risk Scoring Model")

# Scoring table
scoring_rows = [
    ("Critical SAST finding (backdoor, exfiltration)", "Layer 4 \u2014 Diff", "+2.5", RED),
    ("Force push (history rewrite)", "Webhook Event", "+2.0", RED),
    ("CI/CD pipeline tampering (per finding, cap 3.0)", "Layer 4 \u2014 Pipeline", "+1.5", RED),
    ("High SAST finding", "Layer 4 \u2014 Diff", "+1.5", ORANGE),
    ("Author / committer mismatch", "Layer 1 \u2014 Metadata", "+1.5", ORANGE),
    ("Sensitive file touched (per file, cap 3.0)", "Layer 3 \u2014 Files", "+1.5", ORANGE),
    ("Dependency manipulation (supply chain)", "Layer 4 \u2014 Manifests", "+1.5", ORANGE),
    ("Binary file injection (per file, cap 2.0)", "Layer 3 \u2014 Files", "+1.0", ORANGE),
    ("Off-hours commit", "Layer 1 + Baseline", "+1.0", YELLOW),
    ("Large deletion (>500 or >5x baseline)", "Layer 2 + Baseline", "+1.0", YELLOW),
    ("Config / permission weakening (cap 2.0)", "Layer 4 \u2014 Config", "+1.0", YELLOW),
    ("Unsigned commit", "Layer 1 \u2014 Metadata", "+0.5", YELLOW),
    ("Suspicious commit message", "Layer 1 \u2014 Metadata", "+0.5", YELLOW),
]

add_rect(slide, Inches(0.3), Inches(1.3), Inches(8.5), Inches(5.8), fill_color=CARD_BG, border_color=RGBColor(40, 50, 70))
# Header row
y_start = Inches(1.5)
headers = [("Signal", 0.5, 4.0), ("Source", 4.7, 2.3), ("Points", 7.2, 1.2)]
for text, lx, w in headers:
    tb = add_textbox(slide, Inches(lx), y_start, Inches(w), Inches(0.35))
    tf = tb.text_frame
    add_text(tf, text, size=11, bold=True, color=DIM)

for i, (signal, source, pts, color) in enumerate(scoring_rows):
    y = Inches(1.9 + i * 0.38)
    tb = add_textbox(slide, Inches(0.5), y, Inches(4.2), Inches(0.35))
    tf = tb.text_frame
    add_text(tf, signal, size=11, color=LIGHT_GREY)
    tb = add_textbox(slide, Inches(4.7), y, Inches(2.3), Inches(0.35))
    tf = tb.text_frame
    add_text(tf, source, size=11, color=DIM)
    tb = add_textbox(slide, Inches(7.2), y, Inches(1.2), Inches(0.35))
    tf = tb.text_frame
    add_text(tf, pts, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)

# Risk level cards on the right
levels = [
    ("8 \u2013 10", "CRITICAL", "Immediate investigation", RED),
    ("5 \u2013 7.9", "HIGH", "Review within hours", ORANGE),
    ("2.5 \u2013 4.9", "MEDIUM", "Next triage cycle", YELLOW),
    ("0.1 \u2013 2.4", "LOW", "Log for patterns", BLUE),
    ("0", "CLEAN", "No signals detected", GREEN),
]
for i, (score, label, desc, color) in enumerate(levels):
    y = Inches(1.5 + i * 1.1)
    x = Inches(9.2)
    card = add_rect(slide, x, y, Inches(3.8), Inches(0.95), fill_color=CARD_BG, border_color=color)
    tb = add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(1.2), Inches(0.8))
    tf = tb.text_frame
    add_text(tf, score, size=18, bold=True, color=color)
    tb = add_textbox(slide, x + Inches(1.4), y + Inches(0.05), Inches(2.2), Inches(0.8))
    tf = tb.text_frame
    add_text(tf, label, size=12, bold=True, color=color, space_after=Pt(2))
    add_text(tf, desc, size=10, color=GREY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Real-World Scenarios
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "06", "Real-World Insider Threat Scenarios")

scenarios = [
    ("The Departing Engineer", "Employee under notice pushes code at 2 AM that deletes 1,200 lines of auth code and adds a hardcoded admin bypass.",
     "Off-Hours  |  Large Deletion  |  Backdoor SAST  |  Behavioral Anomaly", "Score: 7.0 (High)", RED),
    ("The Credential Harvester", "Developer adds personal AWS key to .env and modifies deploy script to exfiltrate DB connection strings externally.",
     "Sensitive File  |  AWS Key  |  Exfiltration Pattern  |  Connection String", "Score: 8.5 (Critical)", RED),
    ("The Supply Chain Poisoner", "Adds typosquatted npm package with postinstall script, while removing Snyk security scanning from GitHub Actions. Message: 'deps update'.",
     "Dependency Tampering  |  CI/CD Tampering  |  Install Script  |  Suspicious Message", "Score: 5.5 (High)", ORANGE),
    ("The Compromised Account", "Attacker with stolen credentials pushes unsigned commit from different timezone, adding reverse shell in utility file.",
     "Unsigned  |  Off-Hours  |  Reverse Shell SAST  |  Timing Anomaly z=4.8", "Score: 6.0 (High)", ORANGE),
    ("The Quiet Enabler", "Small config changes over weeks: CORS to *, MFA disabled, .env removed from .gitignore, CODEOWNERS emptied.",
     "Config Weakening (x4)  |  Gradual Pattern  |  No SAST Findings", "Score: 4.0/commit (Medium) \u2014 Trend: Increasing", YELLOW),
    ("The Evidence Destroyer", "Force-pushes to main, rewriting 3 weeks of history. Overwritten commits had security audit logs and access control changes.",
     "Force Push  |  History Rewrite  |  Protected Branch  |  Risk Spike 3x", "Score: 5.5 (High)", RED),
]
for i, (title, desc, signals, score, color) in enumerate(scenarios):
    col = i % 3
    row = i // 3
    x = Inches(0.3 + col * 4.3)
    y = Inches(1.4 + row * 3.0)
    card = add_rect(slide, x, y, Inches(4.1), Inches(2.8), fill_color=CARD_BG, border_color=color)
    tb = add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.8), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    add_text(tf, title, size=14, bold=True, color=color, space_after=Pt(4))
    add_text(tf, desc, size=11, color=LIGHT_GREY, space_after=Pt(6))
    add_text(tf, signals, size=10, bold=False, color=DIM, space_after=Pt(4))
    add_text(tf, score, size=11, bold=True, color=color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Competitive Landscape
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "07", "Competitive Landscape \u2014 Market Positioning")

# Left: competitors
add_rect(slide, Inches(0.3), Inches(1.4), Inches(6.3), Inches(5.5), fill_color=CARD_BG, border_color=RGBColor(40, 50, 70))
tb = add_textbox(slide, Inches(0.5), Inches(1.6), Inches(5.9), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True
add_text(tf, "Code-Level Insider Threat Tools", size=16, bold=True, color=WHITE, space_after=Pt(10))

competitors = [
    ("Arnica", "Closest competitor. Scans every commit, has anomalous dev behavior detection. But broader AppSec posture focus, not composite risk scoring."),
    ("Cycode", "Software supply chain security. Code integrity across SDLC, pipeline tampering. Not individual developer risk profiling."),
    ("GitGuardian", "#1 on GitHub Marketplace (500K+ devs). Secrets detection only. No behavioral signals or risk scoring."),
    ("Legit Security", "SDLC posture monitoring. Detects misconfigs in dev infrastructure. Not per-commit analysis."),
]
for name, desc in competitors:
    p = tf.add_paragraph()
    p.space_after = Pt(4)
    r1 = p.add_run()
    r1.text = name
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = ORANGE
    p2 = tf.add_paragraph()
    p2.space_after = Pt(10)
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(11)
    r2.font.color.rgb = GREY

add_text(tf, "Adjacent: DTEX, Proofpoint ITM, Exabeam, Teramind, MS Purview", size=11, color=DIM, space_after=Pt(4))
add_text(tf, "(Endpoint/SIEM-level \u2014 not code-specific)", size=10, color=DIM)

# Right: differentiation
add_rect(slide, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.5), fill_color=CARD_BG, border_color=BLUE)
tb = add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True
add_text(tf, "Our Key Differentiators", size=16, bold=True, color=WHITE, space_after=Pt(10))

diffs = [
    "SAST on commit diffs (58+ insider threat rules) \u2014 nobody else does this",
    "Composite risk scoring (0-10) combining ALL signal types",
    "Per-developer behavioral baselines with z-score anomaly detection",
    "Supply chain signals (typosquatting, vuln pins, custom registries)",
    "CI/CD pipeline tampering detection (11 patterns)",
    "Config weakening detection (9 patterns)",
    "Binary file injection tracking",
    "Developer risk profiles with trend analysis",
    "Commit message anomaly detection",
    "Real-time monitoring without disrupting developer workflow",
]
bullet_list(tf, diffs, size=12, color=LIGHT_GREY, bullet_color=GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Competitive Comparison Matrix
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "08", "Feature Comparison Matrix")

add_rect(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.8), fill_color=CARD_BG, border_color=RGBColor(40, 50, 70))

# Table headers
col_headers = ["Capability", "Arnica", "Cycode", "GitGuardian", "SecureDev AI"]
col_x = [0.5, 5.5, 7.0, 8.5, 10.5]
col_w = [5.0, 1.4, 1.4, 1.8, 2.2]
for i, (text, x, w) in enumerate(zip(col_headers, col_x, col_w)):
    tb = add_textbox(slide, Inches(x), Inches(1.45), Inches(w), Inches(0.4))
    tf = tb.text_frame
    color = GREEN if i == 4 else DIM
    add_text(tf, text, size=11, bold=True, color=color, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

rows = [
    ("SAST on commit diffs", "Partial", "\u2718", "\u2718", "\u2714  58+ IT rules"),
    ("Metadata signals (off-hours, unsigned, mismatch)", "\u2718", "\u2718", "\u2718", "\u2714"),
    ("Sensitive file detection", "\u2718", "\u2718", "Secrets only", "\u2714  Path patterns"),
    ("Binary file injection detection", "\u2718", "\u2718", "\u2718", "\u2714"),
    ("Dependency tampering (typosquatting)", "SCA only", "SCA only", "\u2718", "\u2714  Diff-aware"),
    ("CI/CD pipeline tampering", "\u2718", "Partial", "\u2718", "\u2714  11 patterns"),
    ("Config weakening detection", "\u2718", "\u2718", "\u2718", "\u2714  9 patterns"),
    ("Developer behavioral baseline", "\u2718", "\u2718", "\u2718", "\u2714  Z-score"),
    ("Commit message anomaly", "\u2718", "\u2718", "\u2718", "\u2714"),
    ("Composite risk scoring (0-10)", "Per finding", "\u2718", "\u2718", "\u2714"),
]
for ri, (cap, a, c, g, us) in enumerate(rows):
    y = Inches(1.9 + ri * 0.48)
    vals = [(cap, col_x[0], col_w[0], WHITE, PP_ALIGN.LEFT),
            (a, col_x[1], col_w[1], DIM if a == "\u2718" else YELLOW, PP_ALIGN.CENTER),
            (c, col_x[2], col_w[2], DIM if c == "\u2718" else YELLOW, PP_ALIGN.CENTER),
            (g, col_x[3], col_w[3], DIM if g == "\u2718" else YELLOW, PP_ALIGN.CENTER),
            (us, col_x[4], col_w[4], GREEN, PP_ALIGN.CENTER)]
    for text, x, w, color, align in vals:
        tb = add_textbox(slide, Inches(x), y, Inches(w), Inches(0.4))
        tf = tb.text_frame
        add_text(tf, text, size=11, bold=(text.startswith("\u2714")), color=color, align=align)

# Bottom callout
tb = add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
tf = tb.text_frame
add_text(tf, "Nobody combines SAST findings + metadata signals + behavioral baselines + supply chain signals in a single developer risk profile.", size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Key Differentiators (visual)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
title_slide_header(slide, "09", "Why SecureDev AI")

diff_cards = [
    ("Code-Level Detection", "We analyze actual code diffs, not just metadata or network logs. Every added line is scanned against 58+ insider threat rules.", RED),
    ("Behavioral Intelligence", "Per-developer baselines with rolling 60-commit windows. Z-score anomaly detection personalized to each individual's pattern.", BLUE),
    ("Composite Risk Scoring", "13 signals weighted into a single 0-10 score. Compound signals amplify risk \u2014 no single signal causes a false positive.", ORANGE),
    ("Zero-Friction Deployment", "Connects via GitHub API. No agents to install, no developer workflow changes. Monitoring starts immediately on existing repos.", GREEN),
]
for i, (title, desc, color) in enumerate(diff_cards):
    x = Inches(0.4 + i * 3.2)
    y = Inches(1.8)
    card = add_rect(slide, x, y, Inches(3.0), Inches(4.5), fill_color=CARD_BG, border_color=color)
    tb = add_textbox(slide, x + Inches(0.25), y + Inches(0.5), Inches(2.5), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    add_text(tf, title, size=18, bold=True, color=color, align=PP_ALIGN.CENTER, space_after=Pt(16))
    add_text(tf, desc, size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER, space_after=Pt(8))


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You / Contact
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=RED)

tb = add_textbox(slide, Inches(0), Inches(2.2), W, Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
add_text(tf, "Thank You", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER, space_after=Pt(16))
add_text(tf, "Insider Threat Detection via Commit Intelligence", size=20, color=GREY, align=PP_ALIGN.CENTER, space_after=Pt(24))
add_text(tf, "SecureDev AI", size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/Users/yashwanthgk/appsec-platform/Insider_Threat_Client_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
